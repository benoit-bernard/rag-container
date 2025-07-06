"""
QA Service
Langchain + Ollama + ChromaDB integration
"""
import os
import json
import logging
import time
import gc
from datetime import datetime
from typing import Dict, Any, List, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

class QAService:
    def __init__(self):
        # Basic configuration first
        from app.core.config import config
        
        self.ollama_model = getattr(config, 'OLLAMA_MODEL', 'llama3.2:1b')
        self.ollama_api = getattr(config, 'OLLAMA_BASE_URL', 'http://ollama:11434')
        self.chunk_size = getattr(config, 'CHUNK_SIZE', 1000)
        self.chunk_overlap = getattr(config, 'CHUNK_OVERLAP', 200)
        self.retrieval_k = getattr(config, 'RETRIEVAL_K', 5)
        self.embedding_model = getattr(config, 'EMBEDDING_MODEL', 'sentence-transformers/all-MiniLM-L6-v2')
        
        # Use the volumes paths from docker-compose
        self.persist_dir = str(config.CHROMA_DB_DIR)
        self.documents_dir = str(config.DOCUMENTS_DIR)
        
        # QA Chain (will be initialized when needed)
        self.qa_chain: Optional[Any] = None
        self.last_initialization = None
        self.langchain_available = False
        
        # Check if langchain is available
        try:
            import langchain
            self.langchain_available = True
            logger.info("✅ Langchain available")
        except ImportError:
            logger.warning("⚠️ Langchain not available - QA features will be limited")
        
        logger.info(f"QAService initialized:")
        logger.info(f"  Model: {self.ollama_model}")
        logger.info(f"  API: {self.ollama_api}")
        logger.info(f"  Persist dir: {self.persist_dir}")
        logger.info(f"  Documents dir: {self.documents_dir}")
        logger.info(f"  Langchain available: {self.langchain_available}")
    
    def get_files_registry(self) -> Dict[str, Any]:
        registry = {}
        documents_path = Path(self.documents_dir)
        
        logger.info(f"🔍 Scanning documents directory: {documents_path}")
        logger.info(f"📁 Documents path exists: {documents_path.exists()}")
        
        if not documents_path.exists():
            logger.error(f"❌ Documents directory doesn't exist: {documents_path}")
            return registry
        
        try:
            supported_extensions = [
                '.pdf', '.txt', '.md', 
                '.cs', '.py', ".js", ".cpp", ".c", ".ts",
                ".json", ".xml",
                '.doc', '.docx',            
                '.ppt', '.pptx'     
            ]
            logger.info(f"📂 Supported extensions: {supported_extensions}")

            for filepath in documents_path.rglob("*"):
                if filepath.is_file() and filepath.suffix.lower() in supported_extensions:
                    try:
                        stat = filepath.stat()
                        relative_path = filepath.relative_to(documents_path)
                        registry[str(relative_path)] = {
                            'path': str(filepath),
                            'size': stat.st_size,
                            'mtime': stat.st_mtime,
                            'mtime_readable': datetime.fromtimestamp(stat.st_mtime).isoformat()
                        }
                    except Exception as e:
                        logger.warning(f"Error processing file {filepath}: {e}")
        except Exception as e:
            logger.error(f"Error scanning documents directory: {e}")
        
        logger.info(f"Found {len(registry)} supported files")
        return registry
    
    def initialize_qa_chain(self, force_rebuild: bool = False) -> bool:
        try:
            logger.info("********** 🔄 RAG INITIALIZATION STARTING **********")
            
            if not self.langchain_available:
                logger.warning("********** ❌ LANGCHAIN NOT AVAILABLE **********")
                return False
            
            # Import required modules
            logger.info("********** 📦 IMPORTING LANGCHAIN MODULES **********")
            try:
                from langchain_community.document_loaders import PyPDFLoader, TextLoader
                from langchain.text_splitter import RecursiveCharacterTextSplitter
                from langchain.chains import RetrievalQA
                from langchain_huggingface import HuggingFaceEmbeddings
                from langchain_chroma import Chroma
                from langchain_ollama import OllamaLLM
                import os
                import pickle
                
                logger.info("********** ✅ ALL LANGCHAIN MODULES IMPORTED **********")
            except ImportError as e:
                logger.error(f"********** ❌ LANGCHAIN IMPORT FAILED: {e} **********")
                return False
            
            from app.core.config import config
            read_timeout = getattr(config, 'OLLAMA_READ_TIMEOUT', 300)  # 5 min
            connect_timeout = getattr(config, 'OLLAMA_CONNECT_TIMEOUT', 60)  # 1 min
            
            logger.info(f"********** ⏱️ Using timeouts - Connect: {connect_timeout}s, Read: {read_timeout}s **********")
            
            # Test Ollama connection
            logger.info("********** 🔗 TESTING OLLAMA CONNECTION **********")
            ollama_test = self.test_ollama_connection()
            
            if ollama_test["status"] != "ok":
                logger.error(f"********** ❌ OLLAMA NOT READY: {ollama_test.get('error')} **********")
                return False
            
            logger.info("********** ✅ OLLAMA CONNECTION VERIFIED **********")
            
            logger.info("********** 📁 SETTING UP CHROMADB DIRECTORY **********")
            persist_dir_path = Path(self.persist_dir)
            persist_dir_path.mkdir(parents=True, exist_ok=True)
            
            # ✅ FICHIERS DE CACHE
            cache_metadata_file = persist_dir_path / "documents_cache.json"
            
            logger.info("********** 📄 SCANNING DOCUMENTS **********")
            current_registry = self.get_files_registry()
            if not current_registry:
                logger.warning("********** ⚠️ NO DOCUMENTS FOUND - BASIC LLM MODE **********")
                llm = OllamaLLM(
                    model=self.ollama_model, 
                    base_url=self.ollama_api,
                    timeout=read_timeout,
                    keep_alive=600, 
                    # num_predict=100 
                    )
                self.qa_chain = llm
                self.last_initialization = datetime.now().isoformat()
                logger.info("********** ✅ BASIC LLM INITIALIZED **********")
                return True
            
            logger.info(f"********** 📄 FOUND {len(current_registry)} DOCUMENTS **********")
            
            needs_rebuild = force_rebuild
            cached_registry = {}
            
            if cache_metadata_file.exists() and not force_rebuild:
                try:
                    logger.info("********** 📋 CHECKING DOCUMENT CACHE **********")
                    with open(cache_metadata_file, 'r', encoding='utf-8') as f:
                        cached_registry = json.load(f)
                    
                    if self._compare_registries(current_registry, cached_registry):
                        logger.info("********** ✅ DOCUMENTS UNCHANGED - CHECKING CHROMADB **********")
                        
                        if self._is_chromadb_valid(persist_dir_path):
                            logger.info("********** ✅ CHROMADB VALID - SKIPPING INDEXATION **********")
                            needs_rebuild = False
                        else:
                            logger.info("********** ⚠️ CHROMADB INVALID - REBUILD NEEDED **********")
                            needs_rebuild = True
                    else:
                        logger.info("********** 📝 DOCUMENTS CHANGED - REBUILD NEEDED **********")
                        needs_rebuild = True
                        
                except Exception as e:
                    logger.warning(f"********** ⚠️ CACHE READ ERROR: {e} - REBUILDING **********")
                    needs_rebuild = True
            else:
                logger.info("********** 📝 NO CACHE FOUND - FULL BUILD NEEDED **********")
                needs_rebuild = True
            
            logger.info("********** 🧠 LOADING EMBEDDING MODEL **********")
            embeddings = HuggingFaceEmbeddings(
                model_name=self.embedding_model,
                model_kwargs={'device': 'cpu'}
            )
            
            if needs_rebuild:
                logger.info("********** 🔄 REBUILDING DOCUMENT INDEX **********")
                
                if force_rebuild and persist_dir_path.exists():
                    import shutil
                    shutil.rmtree(persist_dir_path)
                    persist_dir_path.mkdir(parents=True, exist_ok=True)
                    logger.info("********** 🔄 CHROMADB RESET FOR REBUILD **********")
                
                documents = self._load_and_chunk_documents(current_registry)
                
                if not documents:
                    logger.error("********** ❌ NO DOCUMENTS COULD BE LOADED **********")
                    return False
                
                logger.info(f"********** 📊 CREATING VECTORSTORE WITH {len(documents)} CHUNKS **********")
                vectorstore = Chroma.from_documents(
                    documents=documents,
                    embedding=embeddings,
                    persist_directory=str(persist_dir_path)
                )
                
                logger.info("********** 💾 SAVING DOCUMENTS CACHE **********")
                try:
                    with open(cache_metadata_file, 'w', encoding='utf-8') as f:
                        json.dump(current_registry, f, indent=2, ensure_ascii=False)
                    logger.info("********** ✅ CACHE SAVED **********")
                except Exception as e:
                    logger.warning(f"********** ⚠️ CACHE SAVE ERROR: {e} **********")
                
            else:
                logger.info("********** ⚡ LOADING EXISTING CHROMADB **********")
                vectorstore = Chroma(
                    embedding_function=embeddings,
                    persist_directory=str(persist_dir_path)
                )
                logger.info("********** ✅ EXISTING VECTORSTORE LOADED **********")
            
            logger.info("********** 🤖 CREATING OLLAMA LLM **********")
            llm = OllamaLLM(
                model=self.ollama_model,
                base_url=self.ollama_api,
                timeout=read_timeout,      
                keep_alive=600         # 10 min
            )
            
            logger.info("********** 🔗 CREATING RETRIEVAL QA CHAIN **********")
            retriever = vectorstore.as_retriever(
                search_kwargs={"k": self.retrieval_k}
            )
            
            self.qa_chain = RetrievalQA.from_chain_type(
                llm=llm,
                chain_type="stuff",
                retriever=retriever,
                return_source_documents=True
            )
            
            self.last_initialization = datetime.now().isoformat()
            
            status_msg = "REBUILT" if needs_rebuild else "LOADED FROM CACHE"
            logger.info(f"********** ✅ RAG CHAIN {status_msg} SUCCESSFULLY! **********")
            logger.info(f"********** 📄 Documents: {len(current_registry)} **********")
            logger.info(f"********** 🤖 Model: {self.ollama_model} **********")
            logger.info(f"********** 🔍 Retrieval K: {self.retrieval_k} **********")
            logger.info("********** 🎯 RAG SYSTEM READY FOR QUESTIONS **********")
            
            return True
            
        except Exception as e:
            logger.error(f"********** ❌ RAG INITIALIZATION FAILED: {e} **********")
            self.qa_chain = None
            return False
        
    def _compare_registries(self, current: Dict, cached: Dict) -> bool:
        """Compare if document registries are the same"""
        if set(current.keys()) != set(cached.keys()):
            return False
        
        for key in current.keys():
            if key not in cached:
                return False
            if (current[key].get('mtime') != cached[key].get('mtime') or
                current[key].get('size') != cached[key].get('size')):
                return False
        return True
    
    def _is_chromadb_valid(self, persist_dir_path: Path) -> bool:
        """Check if ChromaDB directory is valid"""
        try:
            if not persist_dir_path.exists():
                return False
            
            chroma_files = list(persist_dir_path.glob("*.sqlite*"))
            if not chroma_files:
                return False
            
            from langchain_chroma import Chroma
            from langchain_huggingface import HuggingFaceEmbeddings
            
            embeddings = HuggingFaceEmbeddings(
                model_name=self.embedding_model,
                model_kwargs={'device': 'cpu'}
            )
            
            vectorstore = Chroma(
                embedding_function=embeddings,
                persist_directory=str(persist_dir_path)
            )
            
            # Test
            collection = vectorstore._collection
            count = collection.count()
            logger.info(f"********** 📊 CHROMADB HAS {count} VECTORS **********")
            
            return count > 0
            
        except Exception as e:
            logger.warning(f"********** ⚠️ CHROMADB VALIDATION FAILED: {e} **********")
            return False
    
    def _load_and_chunk_documents(self, registry: Dict) -> List:
        """Load and chunk documents"""
        from langchain_community.document_loaders import (
            PyPDFLoader, PyMuPDFLoader,
            TextLoader,
            UnstructuredWordDocumentLoader,
            UnstructuredPowerPointLoader
        )
        from langchain.text_splitter import RecursiveCharacterTextSplitter
        
        documents = []
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap
        )
        
        processed_count = 0
        for file_info in registry.values():
            try:
                file_path = Path(file_info['path'])
                processed_count += 1
                logger.info(f"********** 📖 [{processed_count}/{len(registry)}] Loading: {file_path.name} **********")
                
                docs = []

                if file_path.suffix.lower() == '.pdf':
                    pdf_loaded = False
                    
                    try:
                        logger.info(f"********** 🔄 Trying PyMuPDFLoader for {file_path.name} **********")
                        loader = PyMuPDFLoader(str(file_path))
                        docs = loader.load()
                        
                        total_content = "".join([doc.page_content for doc in docs])
                        if len(total_content.strip()) > 50:  
                            logger.info(f"********** ✅ PyMuPDFLoader SUCCESS: {len(docs)} pages, {len(total_content)} chars **********")
                            pdf_loaded = True
                        else:
                            logger.warning(f"********** ⚠️ PyMuPDFLoader: Content too short ({len(total_content)} chars) **********")
                            docs = []
                    except Exception as e:
                        logger.warning(f"********** ⚠️ PyMuPDFLoader failed: {e} **********")
                    
                    if not pdf_loaded:
                        try:
                            logger.info(f"********** 🔄 Trying PyPDFLoader for {file_path.name} **********")
                            loader = PyPDFLoader(str(file_path))
                            docs = loader.load()
                            
                            total_content = "".join([doc.page_content for doc in docs])
                            if len(total_content.strip()) > 50:
                                logger.info(f"********** ✅ PyPDFLoader SUCCESS: {len(docs)} pages, {len(total_content)} chars **********")
                                pdf_loaded = True
                            else:
                                logger.warning(f"********** ⚠️ PyPDFLoader: Content too short ({len(total_content)} chars) **********")
                                docs = []
                        except Exception as e:
                            logger.warning(f"********** ⚠️ PyPDFLoader failed: {e} **********")
                    
                    if not pdf_loaded:
                        try:
                            logger.info(f"********** 🔄 Trying manual PDF extraction for {file_path.name} **********")
                            content = self._extract_pdf_manually(file_path)
                            if content and len(content.strip()) > 50:
                                from langchain.schema import Document
                                docs = [Document(
                                    page_content=content,
                                    metadata={"source": str(file_path), "page": 0}
                                )]
                                logger.info(f"********** ✅ Manual extraction SUCCESS: {len(content)} chars **********")
                                pdf_loaded = True
                        except Exception as e:
                            logger.warning(f"********** ⚠️ Manual extraction failed: {e} **********")
                    
                    if not pdf_loaded:
                        logger.error(f"********** ❌ ALL PDF METHODS FAILED for {file_path.name} **********")
                        continue

                elif file_path.suffix.lower() in ['.doc', '.docx']:
                    try:
                        logger.info(f"********** 📄 Loading Word document: {file_path.name} **********")
                        loader = UnstructuredWordDocumentLoader(str(file_path))
                        docs = loader.load()
                        
                        total_content = "".join([doc.page_content for doc in docs])
                        if len(total_content.strip()) > 10:
                            logger.info(f"********** ✅ Word document SUCCESS: {len(total_content)} chars **********")
                        else:
                            logger.warning(f"********** ⚠️ Word document: Content too short ({len(total_content)} chars) **********")
                            continue
                            
                    except Exception as e:
                        logger.warning(f"********** ❌ Word document loading failed: {e} **********")
                        try:
                            content = self._extract_word_manually(file_path)
                            if content and len(content.strip()) > 10:
                                from langchain.schema import Document
                                docs = [Document(
                                    page_content=content,
                                    metadata={"source": str(file_path), "type": "word"}
                                )]
                                logger.info(f"********** ✅ Word manual extraction SUCCESS: {len(content)} chars **********")
                            else:
                                continue
                        except Exception as e2:
                            logger.error(f"********** ❌ Word manual extraction failed: {e2} **********")
                            continue

                elif file_path.suffix.lower() in ['.ppt', '.pptx']:
                    try:
                        logger.info(f"********** 🎯 Loading PowerPoint: {file_path.name} **********")
                        loader = UnstructuredPowerPointLoader(str(file_path))
                        docs = loader.load()
                        
                        total_content = "".join([doc.page_content for doc in docs])
                        if len(total_content.strip()) > 10:
                            logger.info(f"********** ✅ PowerPoint SUCCESS: {len(total_content)} chars **********")
                        else:
                            logger.warning(f"********** ⚠️ PowerPoint: Content too short ({len(total_content)} chars) **********")
                            continue
                            
                    except Exception as e:
                        logger.warning(f"********** ❌ PowerPoint loading failed: {e} **********")
                        try:
                            content = self._extract_powerpoint_manually(file_path)
                            if content and len(content.strip()) > 10:
                                from langchain.schema import Document
                                docs = [Document(
                                    page_content=content,
                                    metadata={"source": str(file_path), "type": "powerpoint"}
                                )]
                                logger.info(f"********** ✅ PowerPoint manual extraction SUCCESS: {len(content)} chars **********")
                            else:
                                continue
                        except Exception as e2:
                            logger.error(f"********** ❌ PowerPoint manual extraction failed: {e2} **********")
                            continue


                elif file_path.suffix.lower() in ['.txt', '.md', '.py', '.cs', ".js", ".cpp", ".c", ".ts", ".json", ".xml"]:
                    try:
                        loader = TextLoader(str(file_path), encoding='utf-8')
                        docs = loader.load()
                        logger.info(f"********** ✅ Text file loaded: {len(docs[0].page_content) if docs else 0} chars **********")
                    except UnicodeDecodeError:
                        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                            try:
                                loader = TextLoader(str(file_path), encoding=encoding)
                                docs = loader.load()
                                logger.info(f"********** ✅ Text file loaded with {encoding}: {len(docs[0].page_content) if docs else 0} chars **********")
                                break
                            except:
                                continue
                        else:
                            logger.error(f"********** ❌ Could not decode text file: {file_path.name} **********")
                            continue
            
                else:
                    logger.warning(f"********** ⏭️ SKIPPING UNSUPPORTED: {file_path.name} **********")
                    continue
                
                if not docs:
                    logger.warning(f"********** ⚠️ NO CONTENT LOADED for {file_path.name} **********")
                    continue
                
                for doc in docs:
                    if hasattr(doc, 'metadata'):
                        if 'source' not in doc.metadata:
                            doc.metadata['source'] = str(file_path)
                        doc.metadata['filename'] = file_path.name
                        doc.metadata['file_path'] = str(file_path)
                        doc.metadata['relative_path'] = str(file_path.relative_to(Path(self.documents_dir)))
                
                total_content = "".join([doc.page_content for doc in docs])
                if len(total_content.strip()) < 10:
                    logger.warning(f"********** ⚠️ CONTENT TOO SHORT for {file_path.name}: {len(total_content)} chars **********")
                    continue

                splits = text_splitter.split_documents(docs)
                documents.extend(splits)
                logger.info(f"********** ✅ {file_path.name}: {len(splits)} chunks created **********")
                
            except Exception as e:
                logger.warning(f"********** ❌ ERROR LOADING {file_path}: {e} **********")
                continue
        
        return documents
    
    def _extract_word_manually(self, file_path: Path) -> str:
        try:
            from docx import Document
            doc = Document(file_path)
            content = ""
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
            return content
        except Exception as e:
            logger.warning(f"Manual Word extraction failed: {e}")
            return ""
    
    def _extract_powerpoint_manually(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = ""
            
            for slide_num, slide in enumerate(prs.slides):
                content += f"\n=== Slide {slide_num + 1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content += shape.text + "\n"
            
            return content
        except Exception as e:
            logger.warning(f"Manual PowerPoint extraction failed: {e}")
            return ""
        
    def test_ollama_connection(self) -> Dict[str, Any]:
        """Test Ollama connection avec retry"""
        import time
        import requests
        
        try:
            health_url = f"{self.ollama_api}/api/tags"
            logger.info(f"🌐 Testing HTTP connection to: {health_url}")
            
            response = requests.get(health_url, timeout=5)
            if response.status_code == 200:
                logger.info("✅ Ollama HTTP endpoint is reachable")
            else:
                logger.warning(f"⚠️ Ollama HTTP returned status: {response.status_code}")
                return {
                    "status": "error",
                    "model": self.ollama_model,
                    "error": f"HTTP status {response.status_code}",
                    "api_url": self.ollama_api
                }
        except requests.exceptions.ConnectionError:
            logger.error("❌ Cannot reach Ollama HTTP endpoint")
            return {
                "status": "error",
                "model": self.ollama_model,
                "error": "Connection refused - Ollama service not ready",
                "api_url": self.ollama_api,
                "suggestion": "Ensure Ollama container is running: docker-compose ps"
            }
        except Exception as e:
            logger.error(f"❌ HTTP test failed: {e}")
            return {
                "status": "error",
                "model": self.ollama_model,
                "error": f"HTTP test failed: {str(e)}",
                "api_url": self.ollama_api
            }
        
        try:
            if not self.langchain_available:
                return {
                    "status": "error",
                    "model": self.ollama_model,
                    "error": "Langchain not available",
                    "api_url": self.ollama_api
                }
            
            from langchain_ollama import OllamaLLM
            
            logger.info(f"🤖 Testing Langchain connection with model: {self.ollama_model}")
            llm = OllamaLLM(
                model=self.ollama_model,
                base_url=self.ollama_api,
                timeout=10  
            )
            
            response = llm.invoke("Hello")
            logger.info(f"✅ Langchain test successful: {response[:50]}...")
            return {
                "status": "ok", 
                "model": self.ollama_model, 
                "response": response[:100],  
                "api_url": self.ollama_api
            }
        except Exception as e:
            logger.error(f"❌ Langchain test failed: {e}")
            return {
                "status": "error", 
                "model": self.ollama_model, 
                "error": str(e),
                "api_url": self.ollama_api,
                "suggestion": f"Model '{self.ollama_model}' might not be downloaded. Try: docker exec rag-ollama ollama pull {self.ollama_model}"
            }
    
    def process_question(self, question: str) -> Dict[str, Any]:
        try:
            from datetime import datetime
            from pathlib import Path
            
            # Ensure QA chain is initialized
            if self.qa_chain is None:
                logger.info("********** QA CHAIN NOT INITIALIZED - ATTEMPTING INITIALIZATION **********")
                if not self.initialize_qa_chain():
                    return {
                        "success": False,
                        "question": question,
                        "answer": "❌ RAG system not ready. Please check Ollama connection and document indexing.",
                        "sources": [],
                        "timestamp": datetime.now().isoformat(),
                        "service_context": {
                            "error": "RAG chain initialization failed",
                            "model": self.ollama_model,
                            "api_url": self.ollama_api
                        }
                    }
            
            logger.info(f"********** ❓ PROCESSING QUESTION: {question} **********")
            
            if hasattr(self.qa_chain, 'invoke') and hasattr(self.qa_chain, 'retriever'):
                logger.info("********** 🔍 USING FULL RAG WITH RETRIEVAL **********")
                
                # Full RAG avec RetrievalQA
                result = self.qa_chain.invoke({"query": question})
                
                answer = result.get("result", "No answer generated")
                source_docs = result.get("source_documents", [])
                
                sources = []
                for i, doc in enumerate(source_docs):
                    try:
                        doc_name = "unknown_document"
                        doc_metadata = {}
                        
                        if hasattr(doc, 'metadata') and doc.metadata:
                            doc_metadata = doc.metadata
                            doc_name = (
                                doc.metadata.get('source') or
                                doc.metadata.get('file') or 
                                doc.metadata.get('filename') or
                                doc.metadata.get('path') or
                                'unknown_document'
                            )
                            
                            if isinstance(doc_name, str) and ('/' in doc_name or '\\' in doc_name):
                                doc_name = Path(doc_name).name
                        
                        content = doc.page_content if hasattr(doc, 'page_content') else str(doc)
                        excerpt = content[:300] + "..." if len(content) > 300 else content
                        
                        score = max(0.9 - (i * 0.15), 0.1)  # De 0.9 à 0.1
                        
                        source_info = {
                            "document": str(doc_name),           # str requis
                            "score": float(score),               # float requis
                            "excerpt": str(excerpt),             # str requis
                            "metadata": doc_metadata             # Optional[Dict[str, Any]]
                        }
                        sources.append(source_info)
                        
                    except Exception as e:
                        logger.warning(f"********** ⚠️ ERROR PROCESSING SOURCE {i}: {e} **********")
                        sources.append({
                            "document": f"document_{i+1}",
                            "score": 0.5,
                            "excerpt": "Content extraction failed",
                            "metadata": {"error": str(e)}
                        })
                
                logger.info(f"********** ✅ RAG ANSWER GENERATED WITH {len(sources)} SOURCES **********")
                
                return {
                    "success": True,
                    "question": question,
                    "answer": answer,
                    "sources": sources, 
                    "timestamp": datetime.now().isoformat(),
                    "service_context": {
                        "model": self.ollama_model,
                        "retrieval_k": len(sources),
                        "processing_mode": "full_rag",
                        "documents_indexed": len(self.get_files_registry()),
                        "api_url": self.ollama_api
                    }
                }
                
            elif hasattr(self.qa_chain, 'invoke'):
                logger.info("********** 🤖 USING SIMPLE LLM (NO RAG) **********")
                
                response = self.qa_chain.invoke(question)
                
                return {
                    "success": True,
                    "question": question,
                    "answer": response,
                    "sources": [],
                    "timestamp": datetime.now().isoformat(),
                    "service_context": {
                        "model": self.ollama_model,
                        "processing_mode": "llm_only",
                        "note": "No RAG - documents not indexed or accessible",
                        "api_url": self.ollama_api
                    }
                }
            else:
                logger.error("********** ❌ QA CHAIN INVALID STATE **********")
                return {
                    "success": False,
                    "question": question,
                    "answer": "❌ QA system is in an invalid state. Please restart the service.",
                    "sources": [],
                    "timestamp": datetime.now().isoformat(),
                    "service_context": {
                        "error": "Invalid QA chain state",
                        "model": self.ollama_model
                    }
                }
            
        except Exception as e:
            logger.error(f"********** ❌ ERROR PROCESSING QUESTION: {e} **********")
            return {
                "success": False,
                "question": question,
                "answer": f"❌ An error occurred while processing your question: {str(e)}",
                "sources": [],
                "timestamp": datetime.now().isoformat(),
                "service_context": {
                    "error": str(e),
                    "model": self.ollama_model,
                    "api_url": self.ollama_api
                }
            }
    
    def get_qa_status(self) -> Dict[str, Any]:
        """Get QA service status"""
        registry = self.get_files_registry()
        
        return {
            "qa_chain_ready": self.qa_chain is not None,
            "langchain_available": self.langchain_available,
            "last_initialization": self.last_initialization,
            "model": self.ollama_model,
            "api_url": self.ollama_api,
            "persist_dir_exists": Path(self.persist_dir).exists(),
            "documents_dir_exists": Path(self.documents_dir).exists(),
            "documents_count": len(registry),
            "config": {
                "chunk_size": self.chunk_size,
                "chunk_overlap": self.chunk_overlap,
                "retrieval_k": self.retrieval_k,
                "embedding_model": self.embedding_model
            }
        }
    
# Global instance
qa_service = QAService()